"""Export the operation-mode logic tree as editable PowerPoint shapes."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "data" / "processed" / "operation_mode_identification"
OUTPUT_PATH = OUTPUT_DIR / "operation_mode_logic_tree_editable.pptx"


COLORS = {
    "start_fill": RGBColor(238, 246, 255),
    "start_line": RGBColor(37, 99, 235),
    "decision_fill": RGBColor(248, 250, 252),
    "decision_line": RGBColor(51, 65, 85),
    "end_fill": RGBColor(236, 253, 245),
    "end_line": RGBColor(5, 150, 105),
    "bad_fill": RGBColor(243, 244, 246),
    "bad_line": RGBColor(107, 114, 128),
    "text": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "arrow": RGBColor(55, 65, 81),
    "label_line": RGBColor(203, 213, 225),
}


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _text(slide, 8.0, 0.28, 8.5, 0.35, "冷站工况识别逻辑树", 24, bold=True)
    _text(slide, 8.0, 0.67, 9.8, 0.25, "先判断制冰，再判断释冰，最后按冷水机组运行组合归类", 11, color=COLORS["muted"])

    nodes = {
        "start": (8.0, 1.05, 2.4, 0.55, "开始\n读取当前时刻数据", "start"),
        "standby": (8.0, 1.95, 3.4, 0.75, "是否停机/无效？\n负荷<500kW且主冷机未运行\n或流量<300m3/h", "decision"),
        "abnormal_top": (11.65, 1.95, 1.7, 0.55, "异常", "bad"),
        "ice_making": (8.0, 3.0, 3.4, 0.86, "是否制冰？\ndual_chiller_on_count >= 1\n且 dual_min_chw_out_temp_c < 0°C", "decision"),
        "ice": (11.65, 3.0, 1.8, 0.58, "制冰", "end"),
        "discharge": (8.0, 4.15, 3.4, 0.72, "是否释冰？\nice_delta_per_step < -50", "decision"),
        "d_base": (4.0, 5.25, 3.2, 0.72, "基载机组是否运行？\nbase_chiller_on_count >= 1", "decision"),
        "m_base": (12.0, 5.25, 3.2, 0.72, "基载机组是否运行？\nbase_chiller_on_count >= 1", "decision"),
        "d_base_dual": (3.0, 6.35, 3.1, 0.72, "双工况机组是否运行？\ndual_chiller_on_count >= 1", "decision"),
        "m_base_dual": (11.0, 6.35, 3.1, 0.72, "双工况机组是否运行？\ndual_chiller_on_count >= 1", "decision"),
        "ice_base": (1.7, 7.7, 2.0, 0.55, "释冰+基载", "end"),
        "ice_base_dual": (4.0, 7.7, 2.15, 0.55, "释冰+基载+双工况", "end"),
        "ice_only": (6.2, 7.7, 1.8, 0.55, "释冰", "end"),
        "base": (9.5, 7.7, 1.7, 0.55, "基载", "end"),
        "base_dual": (11.7, 7.7, 1.95, 0.55, "基载+双工况", "end"),
        "m_abnormal": (14.0, 7.7, 1.6, 0.55, "异常", "bad"),
    }

    edges = [
        ("start", "standby", ""),
        ("standby", "abnormal_top", "Yes"),
        ("standby", "ice_making", "No"),
        ("ice_making", "ice", "Yes"),
        ("ice_making", "discharge", "No"),
        ("discharge", "d_base", "Yes"),
        ("discharge", "m_base", "No"),
        ("d_base", "d_base_dual", "Yes"),
        ("d_base", "ice_only", "No"),
        ("d_base_dual", "ice_base_dual", "Yes"),
        ("d_base_dual", "ice_base", "No"),
        ("m_base", "m_base_dual", "Yes"),
        ("m_base", "m_abnormal", "No"),
        ("m_base_dual", "base_dual", "Yes"),
        ("m_base_dual", "base", "No"),
    ]

    for src, dst, label in edges:
        _connector(slide, nodes[src], nodes[dst], label)

    for node in nodes.values():
        _node(slide, *node)

    _text(
        slide,
        8.0,
        8.72,
        12.0,
        0.25,
        "说明：ice_delta_per_step > 50 不再作为异常判据；若未满足制冰/释冰条件，则按当前冷水机组运行组合归类。",
        10,
        color=COLORS["muted"],
    )

    output_path = OUTPUT_PATH
    try:
        prs.save(output_path)
    except PermissionError:
        output_path = OUTPUT_PATH.with_name(f"{OUTPUT_PATH.stem}_updated{OUTPUT_PATH.suffix}")
        prs.save(output_path)
    print(output_path)


def _node(slide, x: float, y: float, w: float, h: float, text: str, kind: str) -> None:
    fill_key = "start_fill" if kind == "start" else "end_fill" if kind == "end" else "bad_fill" if kind == "bad" else "decision_fill"
    line_key = "start_line" if kind == "start" else "end_line" if kind == "end" else "bad_line" if kind == "bad" else "decision_line"
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - w / 2), Inches(y - h / 2), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill_key]
    shape.line.color.rgb = COLORS[line_key]
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, part in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = part
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(10.5 if kind == "decision" else 12)
        run.font.bold = True
        run.font.color.rgb = COLORS["text"]


def _connector(slide, src, dst, label: str) -> None:
    x0, y0 = _edge_point(src, dst, is_src=True)
    x1, y1 = _edge_point(dst, src, is_src=False)
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    connector.line.color.rgb = COLORS["arrow"]
    connector.line.width = Pt(1.2)
    _add_arrowhead(connector)
    if label:
        lx = (x0 + x1) / 2
        ly = (y0 + y1) / 2
        _label(slide, lx, ly, label)


def _edge_point(node, other, is_src: bool) -> tuple[float, float]:
    x, y, w, h, _, _ = node
    ox, oy, _, _, _, _ = other
    dx = ox - x
    dy = oy - y
    if abs(dx) > abs(dy):
        return (x + (w / 2 if dx > 0 else -w / 2), y)
    return (x, y + (h / 2 if dy > 0 else -h / 2))


def _label(slide, x: float, y: float, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.22), Inches(y - 0.11), Inches(0.44), Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["label_line"]
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = COLORS["text"]


def _text(slide, x: float, y: float, w: float, h: float, text: str, size: int, bold: bool = False, color=None) -> None:
    box = slide.shapes.add_textbox(Inches(x - w / 2), Inches(y - h / 2), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["text"]


def _add_arrowhead(connector) -> None:
    line = connector.line._get_or_add_ln()
    tail = line.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
    if tail is None:
        tail = OxmlElement("a:tailEnd")
        line.append(tail)
    tail.set("type", "triangle")


if __name__ == "__main__":
    main()
