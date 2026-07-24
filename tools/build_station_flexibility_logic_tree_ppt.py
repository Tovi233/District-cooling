from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parents[1] / "outputs" / "冷站侧可调潜力逻辑树_描述版_可编辑.pptx"


COLORS = {
    "title": RGBColor(18, 52, 86),
    "text": RGBColor(24, 35, 48),
    "muted": RGBColor(82, 96, 112),
    "line": RGBColor(42, 83, 145),
    "input_fill": RGBColor(225, 239, 255),
    "calc_fill": RGBColor(239, 246, 255),
    "mode_fill": RGBColor(235, 250, 243),
    "decision_fill": RGBColor(255, 246, 220),
    "output_fill": RGBColor(242, 235, 255),
    "warning_fill": RGBColor(255, 240, 235),
    "white": RGBColor(255, 255, 255),
}


def set_text(shape, title, body="", title_size=9.5, body_size=7.5, color=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(title_size)
    run.font.color.rgb = color or COLORS["title"]

    if body:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = body
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(body_size)
        run.font.color.rgb = COLORS["text"]


def add_box(slide, x, y, w, h, title, body="", fill="calc_fill", title_size=9.5, body_size=7.2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1.2)
    set_text(shape, title, body, title_size, body_size)
    return shape


def add_diamond(slide, x, y, w, h, title, body="", title_size=8.5, body_size=6.8):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["decision_fill"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1.2)
    set_text(shape, title, body, title_size, body_size)
    return shape


def add_arrow(slide, x1, y1, x2, y2, label=None):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = COLORS["line"]
    conn.line.width = Pt(1.1)
    conn.line.end_arrowhead = True
    if label:
        mid_x = (x1 + x2) / 2 - 0.15
        mid_y = (y1 + y2) / 2 - 0.12
        t = slide.shapes.add_textbox(Inches(mid_x), Inches(mid_y), Inches(0.48), Inches(0.22))
        set_text(t, label, "", title_size=6.5, color=COLORS["muted"])
    return conn


def add_title(slide):
    title = slide.shapes.add_textbox(Inches(0.25), Inches(0.12), Inches(15.5), Inches(0.42))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "冷站侧可调潜力快速估算逻辑树（前期阶段）"
    r.font.name = "Microsoft YaHei"
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = COLORS["title"]


def build():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]
    add_title(slide)

    # Row 1: left-to-right preprocessing.
    y1, h1 = 0.78, 1.28
    w1, gap1, x0 = 2.35, 0.22, 0.38
    row1 = []
    row1_data = [
        ("1 输入冷站运行数据", "导入冷站历史运行记录\n作为快速估算的数据基础"),
        ("2 读取当前时刻数据", "读取功率、冷负荷、蓄冰量\n以及当前运行工况"),
        ("3 判断冰量变化", "比较当前与上一时刻蓄冰量\n识别冰量上升或下降"),
        ("4 识别制冰/释冰强度", "由冰量变化判断系统正在\n制冰、释冰或基本不变"),
        ("5 估算可用蓄冰冷量", "扣除安全余量后\n得到可参与响应的蓄冰能力"),
        ("6 估算最大释冰能力", "结合目标响应时长和历史运行能力\n确定可用释冰上限"),
    ]
    for i, (title, body) in enumerate(row1_data):
        row1.append(add_box(slide, x0 + i * (w1 + gap1), y1, w1, h1, title, body, "input_fill"))
    for i in range(len(row1) - 1):
        add_arrow(slide, x0 + (i + 1) * w1 + i * gap1, y1 + h1 / 2, x0 + (i + 1) * (w1 + gap1), y1 + h1 / 2)

    # Row 2: snake turns at right, then operating-mode branch moves right-to-left.
    y2, h2 = 2.52, 1.38
    mode_w, mode_gap = 3.08, 0.22
    modes = [
        ("11 纯释冰 / 异常 / 未定义", "主冷机可削减空间较小\n通常不建议参与削减响应", "warning_fill"),
        ("10 释冰+基载 / 释冰+基载+双工况", "在已有释冰基础上增加释冰强度\n进一步替代机械制冷", "mode_fill"),
        ("9 基载 / 基载+双工况", "优先利用蓄冰替代部分机械制冷\n形成可削减电功率", "mode_fill"),
        ("8 制冰", "通过停止或降低制冰功率\n获得短时削减能力", "mode_fill"),
    ]
    mode_shapes = []
    for i, (title, body, fill) in enumerate(modes):
        mode_shapes.append(add_box(slide, 0.38 + i * (mode_w + mode_gap), y2, mode_w, h2, title, body, fill, 8.1, 6.3))
    decision = add_diamond(slide, 13.40, 2.42, 1.85, 1.60, "7 判断\n当前工况", "按识别规则\n进入对应分支", 8.2, 6.2)
    add_arrow(slide, 14.45, y1 + h1, 14.45, 2.42)
    for s in mode_shapes:
        add_arrow(slide, 13.40, 3.22, s.left.inches + mode_w / 2, y2)

    # Row 3: branches rejoin into power and response calculations.
    y3, h3 = 4.38, 1.22
    r3 = [
        add_box(slide, 0.50, y3, 2.28, h3, "12 冷量转为电功率", "根据制冰或机械制冷效率\n将可转移冷量折算为电功率", "calc_fill", 8.3, 6.8),
        add_box(slide, 3.08, y3, 2.10, h3, "13 得到当前削减能力", "综合停止制冰和蓄冰替代冷机\n得到当前可削减功率", "calc_fill", 8.3, 6.8),
        add_diamond(slide, 5.55, y3 - 0.05, 1.42, 1.35, "14 是否存在\n可削减功率", "", 7.7, 6.3),
        add_box(slide, 7.32, y3, 2.05, h3, "15 否", "当前时段不具备有效削减空间\n响应等级记为不建议响应", "warning_fill", 8.7, 6.4),
        add_box(slide, 9.78, y3, 2.58, h3, "16 是：估算响应时长", "根据蓄冰余量、目标时长\n和当前工况估算可持续时间", "calc_fill", 8.0, 6.2),
        add_box(slide, 12.75, y3, 2.75, h3, "17 估算可转移冷量", "统计响应期间可由蓄冰承担\n或由制冰停止释放的冷量", "calc_fill", 8.3, 6.7),
    ]
    for s in mode_shapes:
        add_arrow(slide, s.left.inches + mode_w / 2, y2 + h2, 1.64, y3)
    for a, b in zip(r3, r3[1:]):
        add_arrow(slide, a.left.inches + a.width.inches, y3 + h3 / 2, b.left.inches, y3 + h3 / 2)
    add_arrow(slide, 6.25, y3 + h3 / 2, 7.32, y3 + h3 / 2, "否")
    add_arrow(slide, 6.95, y3 + h3 / 2, 9.78, y3 + h3 / 2, "是")

    # Row 4: snake continues right-to-left toward final output.
    y4, h4 = 6.35, 1.42
    rebound = add_box(slide, 12.75, y4, 2.75, h4, "18 估算反弹功率", "考虑响应后需要补回的制冰量\n评估后续反弹风险", "calc_fill", 8.4, 6.9)
    level_decision = add_diamond(slide, 10.82, y4 - 0.03, 1.48, 1.48, "19 响应等级\n判断", "", 7.7, 6.2)
    level_ok = add_box(slide, 8.95, y4, 1.50, h4, "20 可响应", "削减能力和持续时间\n均满足响应要求", "mode_fill", 8.2, 6.3)
    level_mid = add_box(slide, 7.15, y4, 1.50, h4, "21 有限响应", "具备一定削减能力\n但持续性或余量有限", "decision_fill", 8.1, 6.2)
    level_no = add_box(slide, 5.35, y4, 1.50, h4, "22 其他", "不建议响应", "warning_fill", 8.2, 6.5)
    output = add_box(
        slide,
        0.38,
        y4,
        4.55,
        h4,
        "23 输出可调潜力指标",
        "输出削减能力、响应时长、可转移冷量\n反弹风险、可用蓄冰能力\n以及最终响应等级",
        "output_fill",
        8.7,
        6.4,
    )
    add_arrow(slide, 14.12, y3 + h3, 14.12, y4)
    add_arrow(slide, 12.75, y4 + h4 / 2, 12.30, y4 + h4 / 2)
    add_arrow(slide, 10.82, y4 + h4 / 2, 10.45, y4 + h4 / 2)
    add_arrow(slide, 10.82, y4 + h4 / 2, level_ok.left.inches + level_ok.width.inches, y4 + h4 / 2)
    add_arrow(slide, 10.82, y4 + h4 / 2, level_mid.left.inches + level_mid.width.inches, y4 + h4 / 2)
    add_arrow(slide, 10.82, y4 + h4 / 2, level_no.left.inches + level_no.width.inches, y4 + h4 / 2)
    add_arrow(slide, 5.35, y4 + h4 / 2, output.left.inches + output.width.inches, y4 + h4 / 2)

    # Small editable parameter note.
    note = slide.shapes.add_textbox(Inches(0.42), Inches(8.22), Inches(15.1), Inches(0.38))
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "说明：本页用于展示冷站侧快速估算的逻辑关系；具体公式、阈值和默认参数保留在项目代码模块中，便于后续统一维护。"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS["muted"]

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build()
